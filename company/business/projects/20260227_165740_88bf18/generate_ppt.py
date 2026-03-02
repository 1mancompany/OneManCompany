from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "如何创建一个成功的 One Man Company (一人公司)"
    subtitle.text = "实事求是，求是创新 | 狼性文化，适者生存"

    slides_data = [
        {
            "title": "什么是 One Man Company？",
            "content": [
                "定义：以一人之力，借助AI和自动化工具，实现传统多部门公司运作的商业模式。",
                "核心理念：极简、高效、自动化、外包。",
                "目标：最大化人效，最小化管理成本。"
            ]
        },
        {
            "title": "One Man Company 的核心支柱",
            "content": [
                "强大的工具链：AI助手、自动化工作流、云服务。",
                "清晰的战略规划：极致聚焦，拒绝贪大求全。",
                "灵活的合作网络：善用外包与自由职业者平台。",
                "极致的自我管理：高度自律与持续迭代的学习能力。"
            ]
        },
        {
            "title": "步骤一：寻找细分市场与精确定位",
            "content": [
                "避开红海：避免与大厂在通用领域正面竞争。",
                "挖掘痛点：寻找长尾需求或垂直领域的具体痛点。",
                "验证MVP：低成本、快速构建最小可行性产品并投入市场验证。",
                "价值主张：确立不可替代的独特价值。"
            ]
        },
        {
            "title": "步骤二：构建自动化工具链与资产库",
            "content": [
                "研发提效：使用大模型辅助编程工具。",
                "运营自动化：利用工具打通各类SaaS，实现工作流闭环。",
                "营销增长：借助AI生成内容，实现自动化分发。",
                "资产沉淀：建立标准操作程序(SOP)，将个人经验转化为系统工具和数字资产。"
            ]
        },
        {
            "title": "步骤三：角色拆解与超级AI协作",
            "content": [
                "角色虚拟化：将CEO、COO、CTO、CMO的角色拆解，分配给对应的AI Agent或自动化流程。",
                "内部沟通机制：建立虚拟会议室、自动化数据报告流，让AI向你汇报。",
                "敏捷迭代：根据业务反馈，随时调整Agent的Prompt和工作流。"
            ]
        },
        {
            "title": "步骤四：建立公司文化与工作准则",
            "content": [
                "价值观：即使是一人公司，也需要明确的价值观（如：实事求是，求是创新）。",
                "工作节奏：设定严格的自我管理标准，避免过度劳累与拖延。",
                "保持狼性：在市场中保持敏锐和竞争力，适者生存。"
            ]
        },
        {
            "title": "挑战与应对策略",
            "content": [
                "挑战：孤独感与驱动力下降。应对：建立外部支持网络，加入创业者社群，定期自我复盘。",
                "挑战：技术或专业瓶颈。应对：善用开源社区，适时引入外部专家或按需外包。",
                "挑战：精力严重分散。应对：严格遵循“重要且紧急”原则，学会对非核心业务说“不”。"
            ]
        },
        {
            "title": "总结与展望",
            "content": [
                "One Man Company 不是规模的限制，而是一种极效的组织形态。",
                "未来趋势：超级个体全面崛起，AI成为核心生产力和基础设施。",
                "立即行动：不要等待完美，从今天的一个简单自动化工作流开始你的旅程！"
            ]
        }
    ]

    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_data["title"]
        tf = body_shape.text_frame
        
        for i, point in enumerate(slide_data["content"]):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point

    prs.save('/Users/yuzhengxu/projects/OneManCompany/company/business/projects/20260227_165740_88bf18/one_man_company_presentation.pptx')
    print("PPT generated successfully at /Users/yuzhengxu/projects/OneManCompany/company/business/projects/20260227_165740_88bf18/one_man_company_presentation.pptx")

if __name__ == '__main__':
    create_presentation()
